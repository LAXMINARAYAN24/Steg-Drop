"""
make_ppt.py — Steg-Drop Presentation Generator
Produces a dark-themed, minimalistic 8-slide PPTX accurately
reflecting the actual codebase: FastAPI + AES-256-GCM + Adaptive LSB.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette ─────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x0D, 0x1A)   # near-black navy
ACCENT    = RGBColor(0x4E, 0xCC, 0xA3)   # teal-green
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE    = RGBColor(0xAA, 0xAA, 0xBB)   # muted text
TAG_BG    = RGBColor(0x1E, 0x1E, 0x35)   # slightly lighter bg for tags

W = Inches(13.333)   # widescreen slide width
H = Inches(7.5)      # widescreen slide height


# ── Helpers ────────────────────────────────────────────────────────────────────

def blank_slide(prs):
    """Add a completely blank slide with the dark background colour."""
    layout = prs.slide_layouts[6]          # blank layout
    slide  = prs.slides.add_slide(layout)
    bg     = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def txt(slide, text, left, top, width, height,
        size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    """Add a text box and return the paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p     = tf.paragraphs[0]
    p.alignment = align
    run   = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    run.font.name      = "Calibri"
    return txBox


def accent_bar(slide, top=Inches(0.55)):
    """Thin horizontal teal rule under the heading."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), top, Inches(12.3), Pt(2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def tag_label(slide, label, left, top, width=Inches(2.8), height=Inches(0.38)):
    """Teal-outlined label chip — used for section tags."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = TAG_BG
    box.line.color.rgb      = ACCENT
    box.line.width          = Pt(1)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top  = tf.margin_bottom = 0
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size  = Pt(11)
    run.font.bold  = True
    run.font.color.rgb = ACCENT
    run.font.name  = "Calibri"


def heading(slide, text):
    """Consistent slide heading at top-left."""
    txt(slide, text,
        left=Inches(0.5), top=Inches(0.15),
        width=Inches(12), height=Inches(0.55),
        size=28, bold=True, color=ACCENT)
    accent_bar(slide)


def bullet_block(slide, items, left, top, width=Inches(12), item_height=Inches(0.42),
                 size=18, color=WHITE, indent=False):
    """Render a list of strings as separate text boxes (avoids pptx placeholder quirks)."""
    y = top
    for item in items:
        c = SUBTLE if indent else color
        s = size - 2 if indent else size
        prefix = "    " if indent else ""
        txt(slide, prefix + item, left, y, width, item_height,
            size=s, color=c)
        y += item_height
    return y


def step_row(slide, number, title, detail, top):
    """Render a numbered step row for workflow slides."""
    # Circle number
    circ = slide.shapes.add_shape(9,  # oval
        Inches(0.5), top + Inches(0.05), Inches(0.42), Inches(0.42))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    tf = circ.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(number)
    r.font.size  = Pt(13)
    r.font.bold  = True
    r.font.color.rgb = BG
    r.font.name  = "Calibri"

    # Title
    txt(slide, title,
        left=Inches(1.1), top=top, width=Inches(3.5), height=Inches(0.5),
        size=16, bold=True, color=WHITE)
    # Detail
    txt(slide, detail,
        left=Inches(4.7), top=top, width=Inches(8.3), height=Inches(0.5),
        size=15, color=SUBTLE)


# ── Slides ─────────────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Title"""
    slide = blank_slide(prs)

    # Top accent line
    bar = slide.shapes.add_shape(1, 0, 0, W, Pt(5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    txt(slide, "STEG-DROP",
        Inches(0.6), Inches(2.2), Inches(11), Inches(1.3),
        size=60, bold=True, color=WHITE)
    txt(slide, "Zero-Knowledge Steganography Service",
        Inches(0.6), Inches(3.45), Inches(11), Inches(0.6),
        size=24, color=ACCENT)
    txt(slide, "Edge-Adaptive LSB  ·  AES-256-GCM  ·  FastAPI",
        Inches(0.6), Inches(4.1), Inches(11), Inches(0.5),
        size=17, color=SUBTLE, italic=True)
    txt(slide, "March 2026",
        Inches(0.6), Inches(5.0), Inches(4), Inches(0.4),
        size=14, color=SUBTLE)

    # Bottom accent line
    bar2 = slide.shapes.add_shape(1, 0, H - Pt(5), W, Pt(5))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = ACCENT
    bar2.line.fill.background()


def slide_intro_problem(prs):
    """Slide 2 — Introduction + Problem (combined, two-column)"""
    slide = blank_slide(prs)
    heading(slide, "Introduction  &  Problem Statement")

    # ── Left column: Introduction ──────────────────────────────────────────────
    txt(slide, "WHAT IS STEG-DROP?",
        Inches(0.5), Inches(0.88), Inches(6.1), Inches(0.35),
        size=12, bold=True, color=ACCENT)

    intro_lines = [
        "Steg-Drop is a zero-knowledge web service that hides",
        "any secret — text, files, keys — inside ordinary images.",
        "",
        "It combines two layers of protection:",
        "  1.  Edge-Adaptive LSB Steganography — embeds bits",
        "       only in high-texture image regions (edge pixels),",
        "       making changes visually and statistically invisible.",
        "  2.  AES-256-GCM Encryption — even if the stego image",
        "       is intercepted, the payload stays unreadable without",
        "       the correct password.",
        "",
        "All processing is in-memory (RAM only). Nothing is ever",
        "written to disk on the server — zero knowledge by design.",
    ]
    y = Inches(1.3)
    for line in intro_lines:
        txt(slide, line, Inches(0.5), y, Inches(6.0), Inches(0.3),
            size=13, color=SUBTLE if line.startswith(" ") else WHITE)
        y += Inches(0.3)

    # Vertical divider
    div = slide.shapes.add_shape(1, Inches(6.75), Inches(0.88), Pt(1.5), Inches(6.3))
    div.fill.solid(); div.fill.fore_color.rgb = TAG_BG
    div.line.fill.background()

    # ── Right column: Problems ─────────────────────────────────────────────────
    txt(slide, "PROBLEMS WE SOLVE",
        Inches(7.0), Inches(0.88), Inches(5.8), Inches(0.35),
        size=12, bold=True, color=ACCENT)

    problems = [
        ("Detectability",
         "Naive sequential LSB is trivially detectable\nby chi-square steganalysis."),
        ("No Confidentiality",
         "Most tools hide data in plaintext — readable\nby anyone who extracts the bits."),
        ("Server Footprint",
         "Online services log, cache, or store uploaded\nimages and payloads on disk."),
        ("Weak Key Derivation",
         "Fixed salts or no key stretching make password-\nbased encryption brute-forceable."),
    ]

    py = Inches(1.3)
    for title, detail in problems:
        # teal bullet dot
        dot = slide.shapes.add_shape(9, Inches(7.0), py + Inches(0.08), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()

        txt(slide, title,
            Inches(7.3), py, Inches(5.5), Inches(0.3),
            size=14, bold=True, color=WHITE)
        txt(slide, detail,
            Inches(7.3), py + Inches(0.3), Inches(5.5), Inches(0.45),
            size=12, color=SUBTLE)
        py += Inches(1.1)


def slide_pipeline(prs):
    """Slide 3 — System Pipeline Diagram"""
    slide = blank_slide(prs)
    heading(slide, "System Pipeline")

    # Helper: draw a rounded-rect node
    def node(left, top, width, height, label, sublabel="", col=ACCENT, text_col=BG):
        shp = slide.shapes.add_shape(
            5,   # rounded rectangle
            left, top, width, height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = col
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.08)
        tf.margin_top  = Inches(0.04)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = text_col
        r.font.name = "Calibri"
        if sublabel:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = sublabel
            r2.font.size = Pt(10)
            r2.font.color.rgb = RGBColor(0x33, 0x33, 0x44) if text_col == BG else SUBTLE
            r2.font.name = "Calibri"

    # Helper: horizontal arrow
    def arrow_h(left, top, width):
        shaft = slide.shapes.add_shape(1, left, top + Inches(0.13), width - Inches(0.12), Pt(2.5))
        shaft.fill.solid(); shaft.fill.fore_color.rgb = ACCENT
        shaft.line.fill.background()
        # arrowhead triangle (right-pointing)
        tri = slide.shapes.add_shape(5, left + width - Inches(0.15), top + Inches(0.05),
                                     Inches(0.15), Inches(0.28))
        tri.fill.solid(); tri.fill.fore_color.rgb = ACCENT
        tri.line.fill.background()

    # Helper: vertical arrow (downward)
    def arrow_v(left, top, height):
        shaft = slide.shapes.add_shape(1, left + Inches(0.13), top, Pt(2.5), height - Inches(0.12))
        shaft.fill.solid(); shaft.fill.fore_color.rgb = RGBColor(0x56, 0xAA, 0xDA)
        shaft.line.fill.background()
        tri = slide.shapes.add_shape(5, left + Inches(0.05), top + height - Inches(0.15),
                                     Inches(0.28), Inches(0.15))
        tri.fill.solid(); tri.fill.fore_color.rgb = RGBColor(0x56, 0xAA, 0xDA)
        tri.line.fill.background()

    # ── Row 1: User → Browser UI → FastAPI ──────────────────────────────────
    NODE_W = Inches(2.1)
    NODE_H = Inches(0.75)
    GAP    = Inches(0.38)
    ROW1_Y = Inches(1.1)

    BLUE  = RGBColor(0x56, 0xAA, 0xDA)
    PURP  = RGBColor(0xDA, 0x70, 0xD6)
    ORG   = RGBColor(0xFF, 0xA5, 0x00)
    RED   = RGBColor(0xFF, 0x6B, 0x6B)

    # Node positions (x)
    x1 = Inches(0.4)
    x2 = x1 + NODE_W + GAP
    x3 = x2 + NODE_W + GAP
    x4 = x3 + NODE_W + GAP
    x5 = x4 + NODE_W + GAP
    x6 = x5 + NODE_W + GAP

    # Row 1 nodes
    node(x1, ROW1_Y, NODE_W, NODE_H, "User",          "Browser",        ACCENT,  BG)
    node(x2, ROW1_Y, NODE_W, NODE_H, "Frontend UI",   "HTML / JS",      BLUE,    WHITE)
    node(x3, ROW1_Y, NODE_W, NODE_H, "FastAPI",        "/encode  /decode",TAG_BG, ACCENT)
    node(x4, ROW1_Y, NODE_W, NODE_H, "Crypto Module",  "PBKDF2 + AES-GCM",PURP, WHITE)
    node(x5, ROW1_Y, NODE_W, NODE_H, "Stego Module",   "Canny + LSB",    ORG,     WHITE)
    node(x6, ROW1_Y, NODE_W, NODE_H, "PNG Stream",     "BytesIO → Client",RED,    WHITE)

    # Row 1 arrows
    for x_start in [x1, x2, x3, x4, x5]:
        arrow_h(x_start + NODE_W, ROW1_Y + Inches(0.24), GAP)

    # ── Labels below each node ───────────────────────────────────────────────
    labels_row1 = [
        (x1, "Upload cover\n+ password"),
        (x2, "POST to\n/encode"),
        (x3, "Orchestrate\nworkflow"),
        (x4, "Encrypt\npayload"),
        (x5, "Embed into\nedge pixels"),
        (x6, "Download\nstego PNG"),
    ]
    for xpos, lbl in labels_row1:
        txt(slide, lbl, xpos, ROW1_Y + NODE_H + Inches(0.06),
            NODE_W, Inches(0.5), size=10, color=SUBTLE, align=PP_ALIGN.CENTER)

    # ── Decode row (below, same pipeline reversed) ───────────────────────────
    ROW2_Y = Inches(3.55)

    node(x6, ROW2_Y, NODE_W, NODE_H, "User",          "Browser",          ACCENT, BG)
    node(x5, ROW2_Y, NODE_W, NODE_H, "Frontend UI",   "HTML / JS",         BLUE,   WHITE)
    node(x4, ROW2_Y, NODE_W, NODE_H, "FastAPI",        "/decode",           TAG_BG, ACCENT)
    node(x3, ROW2_Y, NODE_W, NODE_H, "Stego Module",   "Canny + LSB extract",ORG,   WHITE)
    node(x2, ROW2_Y, NODE_W, NODE_H, "Crypto Module",  "AES-GCM Decrypt",   PURP,  WHITE)
    node(x1, ROW2_Y, NODE_W, NODE_H, "Secret",         "Text / File",        RED,   WHITE)

    # Row 2 arrows (right to left)
    for x_start in [x2, x3, x4, x5]:
        arrow_h(x_start + NODE_W, ROW2_Y + Inches(0.24), GAP)
    # The leftmost arrow (x6→x5)
    arrow_h(x5 + NODE_W, ROW2_Y + Inches(0.24), GAP)

    labels_row2 = [
        (x6, "Upload stego\n+ password"),
        (x5, "POST to\n/decode"),
        (x4, "Route to\ndecoder"),
        (x3, "Extract\nbits"),
        (x2, "Decrypt\n& verify"),
        (x1, "Deliver\nsecret"),
    ]
    for xpos, lbl in labels_row2:
        txt(slide, lbl, xpos, ROW2_Y + NODE_H + Inches(0.06),
            NODE_W, Inches(0.5), size=10, color=SUBTLE, align=PP_ALIGN.CENTER)

    # ── Row labels ──────────────────────────────────────────────────────────
    txt(slide, "ENCODE", Inches(0.4), ROW1_Y - Inches(0.35), Inches(1.5), Inches(0.3),
        size=11, bold=True, color=ACCENT)
    txt(slide, "DECODE", Inches(0.4), ROW2_Y - Inches(0.35), Inches(1.5), Inches(0.3),
        size=11, bold=True, color=BLUE)

    # ── In-Memory note at bottom ─────────────────────────────────────────────
    note_box = slide.shapes.add_shape(1, Inches(0.4), Inches(5.8), Inches(12.5), Inches(0.5))
    note_box.fill.solid(); note_box.fill.fore_color.rgb = TAG_BG
    note_box.line.color.rgb = ACCENT; note_box.line.width = Pt(0.75)
    tf = note_box.text_frame
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "⚡  Zero-Knowledge:  "
    r1.font.bold = True; r1.font.color.rgb = ACCENT; r1.font.size = Pt(13); r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = "All processing happens in server RAM only — no disk writes, no logging, no persistence. Data is garbage-collected after streaming."
    r2.font.color.rgb = SUBTLE; r2.font.size = Pt(12); r2.font.name = "Calibri"


def slide_overview(prs):
    """Slide 3 — System Overview"""
    slide = blank_slide(prs)
    heading(slide, "System Overview")

    layers = [
        ("PRESENTATION LAYER",  "Browser  ·  HTML/JS  ·  Drag-and-drop upload  ·  Stego image download"),
        ("APPLICATION LAYER",   "FastAPI (Uvicorn)  ·  /encode  ·  /decode  ·  Async  ·  Zero disk I/O"),
        ("CRYPTOGRAPHY MODULE", "PBKDF2-HMAC-SHA256  ·  AES-256-GCM  ·  Random salt  ·  Payload serialization"),
        ("STEGANOGRAPHY MODULE","Canny Edge Detection  ·  Pixel edge map  ·  Capacity check  ·  LSB embedding"),
        ("IN-MEMORY STREAM",    "All processing in RAM  ·  PNG streamed to client  ·  Data GC'd after response"),
    ]

    accent_colors = [
        RGBColor(0x4E, 0xCC, 0xA3),
        RGBColor(0x56, 0xAA, 0xDA),
        RGBColor(0xDA, 0x70, 0xD6),
        RGBColor(0xFF, 0xA5, 0x00),
        RGBColor(0xFF, 0x6B, 0x6B),
    ]

    y = Inches(0.9)
    for i, (layer, detail) in enumerate(layers):
        col = accent_colors[i]
        # Side accent stripe
        stripe = slide.shapes.add_shape(1, Inches(0.5), y, Pt(4), Inches(0.75))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = col
        stripe.line.fill.background()

        txt(slide, layer,
            Inches(0.75), y, Inches(4.5), Inches(0.4),
            size=13, bold=True, color=col)
        txt(slide, detail,
            Inches(0.75), y + Inches(0.32), Inches(12), Inches(0.45),
            size=13, color=SUBTLE)

        if i < len(layers) - 1:
            arr = slide.shapes.add_shape(1, Inches(0.62), y + Inches(0.78), Pt(2), Inches(0.12))
            arr.fill.solid(); arr.fill.fore_color.rgb = SUBTLE
            arr.line.fill.background()

        y += Inches(1.0)


def slide_crypto(prs):
    """Slide 4 — Cryptography Module"""
    slide = blank_slide(prs)
    heading(slide, "Cryptography Module")

    sections = [
        ("KEY DERIVATION",
         ["Algorithm:   PBKDF2-HMAC-SHA256",
          "Iterations:   100,000  (OWASP minimum recommendation)",
          "Salt:          16-byte cryptographically random (per encode)",
          "Output:        32-byte (256-bit) AES key"]),
        ("ENCRYPTION  —  AES-256-GCM",
         ["Mode:          GCM (Galois/Counter Mode) — authenticated encryption",
          "Nonce:         12-byte random per encryption operation",
          "Tag:           16-byte authentication tag — detects any tampering",
          "Wire format:   salt (16B) | nonce (12B) | ciphertext | tag (16B)"]),
        ("PAYLOAD SERIALIZATION",
         ["Format:        type_flag (1B) | filename_len (2B) | filename | raw data",
          "Supports:      text messages and arbitrary binary files"]),
    ]

    y = Inches(0.9)
    for title, bullets in sections:
        txt(slide, title, Inches(0.5), y, Inches(12), Inches(0.35),
            size=13, bold=True, color=ACCENT)
        y += Inches(0.35)
        for b in bullets:
            txt(slide, "  " + b, Inches(0.5), y, Inches(12.3), Inches(0.33),
                size=14, color=SUBTLE)
            y += Inches(0.33)
        y += Inches(0.18)


def slide_stego(prs):
    """Slide 5 — Steganography Module"""
    slide = blank_slide(prs)
    heading(slide, "Steganography Module  —  Adaptive LSB")

    txt(slide,
        "Data is embedded exclusively in high-texture edge pixels, making modifications statistically indistinguishable from natural image noise.",
        Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.55),
        size=16, color=SUBTLE)

    steps = [
        ("1. Grey-scale conversion", "RGB cover image → single-channel grey for edge analysis."),
        ("2. LSB zeroing",           "Zero the LSB of all pixels so the Canny mask is identical before & after embedding."),
        ("3. Canny edge detection",  "OpenCV Canny with auto-thresholding (0.5× / 1.5× median). Finds complex, busy regions."),
        ("4. Mask dilation",         "5×5 elliptical kernel, 2 iterations — expands the embeddable zone around each edge."),
        ("5. Capacity verification", "Count edge pixels × 3 channels × 1 bit = max bits. Reserve 4 bytes for length header."),
        ("6. LSB substitution",      "Walk edge positions sequentially, replace LSB of R, G, B with payload bits."),
        ("7. Lossless PNG export",   "Save stego array as PNG (not JPEG) — lossless format guarantees no bit corruption."),
    ]

    y = Inches(1.5)
    for title, detail in steps:
        circle = slide.shapes.add_shape(9,
            Inches(0.5), y + Inches(0.05), Inches(0.32), Inches(0.32))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()

        txt(slide, title, Inches(0.95), y, Inches(4.0), Inches(0.42),
            size=14, bold=True, color=WHITE)
        txt(slide, detail, Inches(5.05), y, Inches(8.0), Inches(0.42),
            size=13, color=SUBTLE)
        y += Inches(0.7)


def slide_encode(prs):
    """Slide 6 — Encode Workflow"""
    slide = blank_slide(prs)
    heading(slide, "Encoding Workflow")

    tag_label(slide, "POST /encode", Inches(9.8), Inches(0.12), width=Inches(2.85))

    steps_data = [
        ("Serialize Payload",   "Bundle secret (text/file) with type flag, filename length, filename, raw bytes."),
        ("Generate Salt",       "os.urandom(16) — fresh random salt for every encode operation."),
        ("Derive AES Key",      "PBKDF2-HMAC-SHA256(password, salt, 100 000 iters) → 256-bit key."),
        ("AES-256-GCM Encrypt", "Encrypt serialized payload; prepend 16-byte salt to wire payload."),
        ("Canny Edge Mask",     "Build boolean pixel mask from Canny filter on the cover image."),
        ("Capacity Check",      "Assert encrypted_size ≤ available edge-pixel bits, raise HTTP 400 if not."),
        ("LSB Embedding",       "Write each encrypted bit into R/G/B LSBs of each edge pixel, in order."),
        ("Stream PNG",          "Save stego NumPy array → PNG (BytesIO) → StreamingResponse. Zero disk writes."),
    ]

    y = Inches(0.88)
    for i, (title, detail) in enumerate(steps_data, 1):
        step_row(slide, i, title, detail, y)
        y += Inches(0.76)


def slide_decode(prs):
    """Slide 7 — Decode Workflow"""
    slide = blank_slide(prs)
    heading(slide, "Decoding Workflow")

    tag_label(slide, "POST /decode", Inches(9.8), Inches(0.12), width=Inches(2.85))

    steps_data = [
        ("Receive Stego PNG",    "Accept stego image upload + master password from client."),
        ("Canny Mask (same)",    "Re-run identical Canny filter — deterministically reproduces the original edge map."),
        ("Read Length Header",   "Extract first 32 bits from LSBs of edge pixels → 4-byte big-endian payload length."),
        ("Extract Ciphertext",   "Pull exactly length × 8 + 32 bits from edge pixels; pack to bytes."),
        ("Separate Salt",        "First 16 bytes of extracted data = random salt used during encoding."),
        ("Re-derive AES Key",    "PBKDF2(password, salt, 100 000 iters) → reconstructs the same 256-bit key."),
        ("AES-256-GCM Decrypt",  "Decrypt and verify GCM tag — raises ValueError on wrong password or tampered data."),
        ("Deserialize & Deliver","Unpack type flag + filename + data; return JSON (text) or StreamingResponse (file)."),
    ]

    y = Inches(0.88)
    for i, (title, detail) in enumerate(steps_data, 1):
        step_row(slide, i, title, detail, y)
        y += Inches(0.76)


def slide_stack(prs):
    """Slide 8 — Tech Stack & Future Work"""
    slide = blank_slide(prs)
    heading(slide, "Tech Stack  ·  Future Work")

    # ----- Tech Stack (left column) -----
    txt(slide, "TECH STACK", Inches(0.5), Inches(0.9), Inches(6), Inches(0.35),
        size=13, bold=True, color=ACCENT)

    stack = [
        ("FastAPI + Uvicorn",      "Async Python web framework / ASGI server"),
        ("Pillow",                 "Image open / convert / save (PNG lossless)"),
        ("OpenCV (headless)",      "Canny edge detection on NumPy arrays"),
        ("NumPy",                  "Bit manipulation — packbits / unpackbits"),
        ("PyCryptodome",           "AES-256-GCM via Crypto.Cipher.AES"),
        ("hashlib (stdlib)",       "PBKDF2-HMAC-SHA256 key derivation"),
        ("python-multipart",       "FastAPI form / file upload parsing"),
    ]

    y = Inches(1.3)
    for lib, purpose in stack:
        txt(slide, lib,     Inches(0.5), y, Inches(3.0), Inches(0.35),
            size=14, bold=True, color=WHITE)
        txt(slide, purpose, Inches(3.6), y, Inches(3.0), Inches(0.35),
            size=13, color=SUBTLE)
        y += Inches(0.55)

    # Divider
    div = slide.shapes.add_shape(1, Inches(6.9), Inches(0.9), Pt(1.5), Inches(6.3))
    div.fill.solid(); div.fill.fore_color.rgb = TAG_BG
    div.line.fill.background()

    # ----- Future Work (right column) -----
    txt(slide, "FUTURE WORK", Inches(7.1), Inches(0.9), Inches(5.7), Inches(0.35),
        size=13, bold=True, color=ACCENT)

    roadmap = [
        ("🛡  Multi-bit LSB",       "Embed 2–3 bits per channel to increase capacity without visible artefacts."),
        ("🌐  WASM Frontend",        "Move encoding entirely client-side via WebAssembly — server never sees plaintext."),
        ("🎥  Video Steganography",  "Embed data in H.264 motion-vector residuals for video cover media."),
        ("🔊  Audio Steganography",  "LSB embedding in PCM WAV files as an alternate carrier channel."),
        ("📊  Steganalysis Evasion", "Adaptive embedding density based on local image entropy score."),
        ("🔑  Hardware Keys",        "FIDO2/YubiKey support for password-free key derivation."),
    ]

    y = Inches(1.3)
    for icon_title, detail in roadmap:
        txt(slide, icon_title, Inches(7.1), y, Inches(5.7), Inches(0.32),
            size=14, bold=True, color=WHITE)
        txt(slide, detail, Inches(7.1), y + Inches(0.3), Inches(5.7), Inches(0.3),
            size=12, color=SUBTLE)
        y += Inches(0.82)


# ── Main ────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_intro_problem(prs)
    slide_pipeline(prs)
    slide_overview(prs)
    slide_crypto(prs)
    slide_stego(prs)
    slide_encode(prs)
    slide_decode(prs)
    slide_stack(prs)

    out = r"C:\Users\sahul\Desktop\steg-drop\StegDrop_Presentation.pptx"
    prs.save(out)
    print(f"Done. Saved → {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
